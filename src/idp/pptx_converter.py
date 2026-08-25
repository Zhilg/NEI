"""Convert PPTX files to Markdown using python-pptx."""

from __future__ import annotations

from pathlib import Path

from pptx import Presentation


def convert_pptx_to_markdown(pptx_path: Path) -> str:
    prs = Presentation(str(pptx_path))
    parts: list[str] = []
    for i, slide in enumerate(prs.slides, 1):
        texts: list[str] = []
        for shape in slide.shapes:
            try:
                if hasattr(shape, "text") and shape.text.strip():
                    texts.append(shape.text.strip())
            except Exception:
                continue
        if texts:
            parts.append(f"<!-- slide {i} -->\n" + "\n".join(texts))
    return "\n\n".join(parts)
